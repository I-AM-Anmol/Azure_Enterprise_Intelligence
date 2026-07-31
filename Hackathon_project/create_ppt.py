"""
Generate Professional Hackathon Presentation
Clinical No-Show Prediction — Infinity Nexus Team
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Color Palette (matching dashboard theme)
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x00, 0xB8, 0xA9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0x94, 0xA3, 0xB8)
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)
CARD_BG = RGBColor(0x1E, 0x30, 0x50)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
RED = RGBColor(0xC6, 0x28, 0x28)
OFF_WHITE = RGBColor(0xFA, 0xFB, 0xFC)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_text(slide, text, left, top, width, height, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def add_bullet_text(slide, bullets, left, top, width, height, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_before = Pt(6)
        p.level = 0
    return txBox


def add_shape_box(slide, left, top, width, height, text, fill_color, text_color=WHITE, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = TEAL
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, start_left, start_top, end_left, end_top):
    connector = slide.shapes.add_connector(
        1, start_left, start_top, end_left, end_top
    )
    connector.line.color.rgb = TEAL
    connector.line.width = Pt(2)
    return connector


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: TITLE
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "🏥", Inches(5.8), Inches(1.2), Inches(2), Inches(1),
                   font_size=60, align=PP_ALIGN.CENTER)
    add_title_text(slide, "Clinical No-Show Prediction",
                   Inches(1.5), Inches(2.5), Inches(10), Inches(1),
                   font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_title_text(slide, "AI-Powered Patient Risk Stratification & Automated Intervention",
                   Inches(1.5), Inches(3.5), Inches(10), Inches(0.8),
                   font_size=20, bold=False, color=TEAL, align=PP_ALIGN.CENTER)
    add_title_text(slide, "Infinity Nexus Team  |  Hackathon 2026",
                   Inches(1.5), Inches(5.0), Inches(10), Inches(0.6),
                   font_size=18, bold=False, color=LIGHT_GREY, align=PP_ALIGN.CENTER)
    add_title_text(slide, "Powered by: Power BI  •  XGBoost  •  Streamlit  •  Power Automate  •  Groq AI",
                   Inches(1.5), Inches(5.7), Inches(10), Inches(0.6),
                   font_size=14, bold=False, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: PROBLEM STATEMENT
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "The Problem", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    problems = [
        "📉  Patient no-shows cost U.S. healthcare $150B+ annually",
        "🏥  Average clinic loses $200-$500 per missed appointment",
        "📅  No-shows disrupt scheduling, waste staff time, delay care",
        "❌  Traditional reminder systems are one-size-fits-all",
        "⚠️  High-risk patients get the same reminder as low-risk ones",
    ]
    add_bullet_text(slide, problems, Inches(0.8), Inches(1.5), Inches(7), Inches(4), font_size=18)

    # Key stat box
    add_shape_box(slide, Inches(8.5), Inches(2.0), Inches(4), Inches(2.5),
                  "Industry Average\nNo-Show Rate\n\n23-34%",
                  CARD_BG, WHITE, font_size=18)

    add_title_text(slide, "What if we could predict WHO will no-show and intervene BEFORE it happens?",
                   Inches(0.8), Inches(6.0), Inches(11), Inches(0.8),
                   font_size=16, bold=True, color=TEAL, align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: SOLUTION OVERVIEW
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Our Solution", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    add_title_text(slide, "End-to-end ML pipeline that predicts no-show risk and triggers\n"
                   "targeted interventions — all from a single operational dashboard.",
                   Inches(0.8), Inches(1.3), Inches(11), Inches(1.2),
                   font_size=18, bold=False, color=WHITE)

    # Three pillars
    add_shape_box(slide, Inches(0.5), Inches(3.0), Inches(3.8), Inches(3.0),
                  "PREDICT\n━━━━━━\nXGBoost ML model\ntrained on 4,800+\nhistorical visits\n\nAUC > 0.55",
                  CARD_BG, WHITE, font_size=14)

    add_shape_box(slide, Inches(4.8), Inches(3.0), Inches(3.8), Inches(3.0),
                  "STRATIFY\n━━━━━━\n3 Risk Tiers:\n🟢 Low (<30%)\n🟡 Medium (30-60%)\n🔴 High (>60%)",
                  CARD_BG, WHITE, font_size=14)

    add_shape_box(slide, Inches(9.1), Inches(3.0), Inches(3.8), Inches(3.0),
                  "INTERVENE\n━━━━━━\nPower Automate\nemail reminders\nto high-risk patients\nwith one click",
                  CARD_BG, WHITE, font_size=14)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: ARCHITECTURE DIAGRAM
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Solution Architecture", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    # Row 1: Data Sources
    add_shape_box(slide, Inches(0.5), Inches(1.8), Inches(2.5), Inches(1.2),
                  "📊 Power BI\nSemantic Model", RGBColor(0x24, 0x33, 0x52), WHITE, font_size=12)

    add_shape_box(slide, Inches(3.5), Inches(1.8), Inches(2.5), Inches(1.2),
                  "🔐 Azure AD\nService Principal", RGBColor(0x24, 0x33, 0x52), WHITE, font_size=12)

    # Row 2: Processing
    add_shape_box(slide, Inches(0.5), Inches(3.8), Inches(2.5), Inches(1.2),
                  "📝 DAX Queries\n(REST API)", RGBColor(0x24, 0x33, 0x52), WHITE, font_size=12)

    add_shape_box(slide, Inches(3.5), Inches(3.8), Inches(2.5), Inches(1.2),
                  "🤖 XGBoost\nML Model", TEAL, WHITE, font_size=12)

    add_shape_box(slide, Inches(6.5), Inches(3.8), Inches(2.5), Inches(1.2),
                  "📊 Risk\nStratification", RGBColor(0x24, 0x33, 0x52), WHITE, font_size=12)

    # Row 3: Output
    add_shape_box(slide, Inches(0.5), Inches(5.8), Inches(2.5), Inches(1.2),
                  "📱 Streamlit\nDashboard", RGBColor(0x24, 0x33, 0x52), WHITE, font_size=12)

    add_shape_box(slide, Inches(3.5), Inches(5.8), Inches(2.5), Inches(1.2),
                  "🤖 AI Chatbot\n(Groq/Llama 3.1)", RGBColor(0x24, 0x33, 0x52), WHITE, font_size=12)

    add_shape_box(slide, Inches(6.5), Inches(5.8), Inches(2.5), Inches(1.2),
                  "📧 Power Automate\nEmail Reminders", RGBColor(0x24, 0x33, 0x52), WHITE, font_size=12)

    # Flow arrows (simplified as text)
    add_title_text(slide, "→", Inches(3.0), Inches(2.0), Inches(0.5), Inches(0.8),
                   font_size=24, color=TEAL, align=PP_ALIGN.CENTER)
    add_title_text(slide, "↓", Inches(1.5), Inches(3.1), Inches(0.5), Inches(0.6),
                   font_size=24, color=TEAL, align=PP_ALIGN.CENTER)
    add_title_text(slide, "→", Inches(3.0), Inches(4.0), Inches(0.5), Inches(0.8),
                   font_size=24, color=TEAL, align=PP_ALIGN.CENTER)
    add_title_text(slide, "→", Inches(6.0), Inches(4.0), Inches(0.5), Inches(0.8),
                   font_size=24, color=TEAL, align=PP_ALIGN.CENTER)
    add_title_text(slide, "↓", Inches(1.5), Inches(5.1), Inches(0.5), Inches(0.6),
                   font_size=24, color=TEAL, align=PP_ALIGN.CENTER)
    add_title_text(slide, "↓", Inches(4.5), Inches(5.1), Inches(0.5), Inches(0.6),
                   font_size=24, color=TEAL, align=PP_ALIGN.CENTER)
    add_title_text(slide, "↓", Inches(7.5), Inches(5.1), Inches(0.5), Inches(0.6),
                   font_size=24, color=TEAL, align=PP_ALIGN.CENTER)

    # Legend
    add_title_text(slide, "Data Flow: Power BI → DAX API → Feature Engineering → XGBoost → Risk Scores → Dashboard + Automation",
                   Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                   font_size=11, bold=False, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: DATA PIPELINE
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Data Pipeline", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    items = [
        "📊  Source: Power BI Semantic Model (workspace: Clinical No-Show Prediction)",
        "🔐  Auth: Azure AD Service Principal (client_credentials OAuth2 flow)",
        "📝  Queries: DAX via Power BI REST API (executeQueries endpoint)",
        "📋  Tables: staging_appointment (14 cols) + staging_patient (16 cols)",
        "🔗  Join: Left merge on patient_id → 4,800+ enriched records",
        "⚙️  Feature Engineering: lead_days, day_of_week, hour, past_noshow_ratio, etc.",
        "⏱️  Caching: st.cache_data with 1-hour TTL for performance",
    ]
    add_bullet_text(slide, items, Inches(0.5), Inches(1.5), Inches(12), Inches(5), font_size=17)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: ML MODEL
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "ML Model — XGBoost Classifier", Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
                   font_size=32, color=TEAL)

    # Model details
    model_info = [
        "Algorithm: XGBClassifier (Extreme Gradient Boosting)",
        "Training Data: 4,800+ historical in-person appointments",
        "Train/Test Split: 80/20 with stratified sampling",
        "Evaluation: ROC-AUC Score",
        "Class Balancing: scale_pos_weight (auto-computed)",
    ]
    add_bullet_text(slide, model_info, Inches(0.5), Inches(1.4), Inches(6), Inches(3), font_size=16)

    # Hyperparameters box
    add_shape_box(slide, Inches(7.5), Inches(1.4), Inches(5.2), Inches(3.2),
                  "Hyperparameters\n━━━━━━━━━━━━\nn_estimators: 500\nmax_depth: 6\n"
                  "learning_rate: 0.05\nsubsample: 0.85\ncolsample_bytree: 0.7\n"
                  "min_child_weight: 3\nreg_alpha: 0.1 | reg_lambda: 1.5",
                  CARD_BG, WHITE, font_size=13)

    # Risk tiers
    add_title_text(slide, "Risk Stratification Thresholds:", Inches(0.5), Inches(4.8), Inches(6), Inches(0.5),
                   font_size=16, bold=True, color=WHITE)

    add_shape_box(slide, Inches(0.5), Inches(5.5), Inches(3.5), Inches(1.2),
                  "🟢 Low Risk\n< 30% probability", GREEN, WHITE, font_size=14)
    add_shape_box(slide, Inches(4.5), Inches(5.5), Inches(3.5), Inches(1.2),
                  "🟡 Medium Risk\n30% – 60%", ORANGE, WHITE, font_size=14)
    add_shape_box(slide, Inches(8.5), Inches(5.5), Inches(3.5), Inches(1.2),
                  "🔴 High Risk\n> 60% probability", RED, WHITE, font_size=14)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: FEATURE IMPORTANCE
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Top Predictive Features", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    features_ranked = [
        ("past_noshow_ratio", "Patient's historical no-show percentage", "★★★★★"),
        ("risk_combo", "past_noshow × lead_days interaction", "★★★★☆"),
        ("lead_days", "Days between booking and appointment", "★★★★☆"),
        ("distance_to_clinic_miles", "How far patient lives from clinic", "★★★★☆"),
        ("patient_age", "Patient age at appointment", "★★★☆☆"),
        ("total_past_appts", "Patient engagement history", "★★★☆☆"),
        ("hour", "Time of day (early = higher risk)", "★★☆☆☆"),
        ("day_of_week", "Monday/Friday higher risk", "★★☆☆☆"),
        ("no_sms", "Not enrolled in SMS reminders", "★★☆☆☆"),
        ("insurance_encoded", "Self-pay = higher risk", "★☆☆☆☆"),
    ]

    y_pos = Inches(1.4)
    for feat, desc, stars in features_ranked:
        add_title_text(slide, f"{feat}", Inches(0.5), y_pos, Inches(3.5), Inches(0.4),
                       font_size=13, bold=True, color=WHITE)
        add_title_text(slide, desc, Inches(4.2), y_pos, Inches(5.5), Inches(0.4),
                       font_size=12, bold=False, color=LIGHT_GREY)
        add_title_text(slide, stars, Inches(10.5), y_pos, Inches(2.5), Inches(0.4),
                       font_size=13, bold=False, color=TEAL)
        y_pos += Inches(0.55)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8: DASHBOARD FEATURES
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Streamlit Dashboard Features", Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
                   font_size=32, color=TEAL)

    features_list = [
        "📊  Real-time KPI Cards — High/Medium/Low risk counts + revenue at risk",
        "🎯  Risk Stratification — Color-coded cards with recommended actions",
        "📅  Weekly & Daily Views — Appointment volume breakdown by risk tier",
        "📋  Scored Table — All upcoming appointments with sortable risk scores",
        "💰  Cost Impact Analysis — Revenue loss estimation & recoverable savings",
        "📈  Historical Trends — Monthly no-show rate trend line",
        "🔍  Risk Factor Breakdown — Insurance, day-of-week, category analysis",
        "🎛️  Sidebar Filters — Date range, risk tier, appointment category",
        "🔄  Live Refresh — One-click data refresh from Power BI model",
    ]
    add_bullet_text(slide, features_list, Inches(0.5), Inches(1.4), Inches(12), Inches(5.5), font_size=16)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 9: AI CHATBOT
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "AI-Powered Chatbot", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    chatbot_info = [
        "🤖  Model: Llama 3.1 8B Instant (via Groq API — free tier)",
        "💬  Location: Sidebar expander for quick access",
        "📊  Context: Fed with real-time data summary + top 50 patient details",
        "🔍  Capabilities:",
        "      • \"How many patients are at high risk this week?\"",
        "      • \"Name the patients with highest no-show probability\"",
        "      • \"What are the top reasons for no-shows?\"",
        "      • \"Which insurance type has most no-shows?\"",
        "⚡  Response time: < 2 seconds (Groq inference)",
        "🗑️  Clear chat history with one click",
    ]
    add_bullet_text(slide, chatbot_info, Inches(0.5), Inches(1.4), Inches(8), Inches(5.5), font_size=16)

    # Chatbot visual
    add_shape_box(slide, Inches(9.0), Inches(1.8), Inches(3.8), Inches(4.5),
                  "🤖 No-Show\nIntelligence\nAssistant\n━━━━━━━━\n\n"
                  "User: How many\nare high risk?\n\n"
                  "Bot: There are 5\nhigh-risk patients\nthis week...",
                  CARD_BG, WHITE, font_size=12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 10: POWER AUTOMATE
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Power Automate Integration", Inches(0.5), Inches(0.3), Inches(10), Inches(0.8),
                   font_size=32, color=TEAL)

    add_title_text(slide, "One-click email reminders to high-risk patients",
                   Inches(0.5), Inches(1.2), Inches(10), Inches(0.6),
                   font_size=18, bold=False, color=LIGHT_GREY)

    # Flow diagram
    add_shape_box(slide, Inches(0.3), Inches(2.5), Inches(2.4), Inches(1.5),
                  "📧 Click\n\"Send Reminders\"\nButton", TEAL, WHITE, font_size=12)

    add_title_text(slide, "→", Inches(2.7), Inches(2.9), Inches(0.5), Inches(0.8),
                   font_size=28, color=TEAL, align=PP_ALIGN.CENTER)

    add_shape_box(slide, Inches(3.3), Inches(2.5), Inches(2.4), Inches(1.5),
                  "📤 POST JSON\nto Power Automate\nWebhook", CARD_BG, WHITE, font_size=12)

    add_title_text(slide, "→", Inches(5.7), Inches(2.9), Inches(0.5), Inches(0.8),
                   font_size=28, color=TEAL, align=PP_ALIGN.CENTER)

    add_shape_box(slide, Inches(6.3), Inches(2.5), Inches(2.4), Inches(1.5),
                  "⚙️ Parse JSON\n& Loop Each\nPatient", CARD_BG, WHITE, font_size=12)

    add_title_text(slide, "→", Inches(8.7), Inches(2.9), Inches(0.5), Inches(0.8),
                   font_size=28, color=TEAL, align=PP_ALIGN.CENTER)

    add_shape_box(slide, Inches(9.3), Inches(2.5), Inches(2.4), Inches(1.5),
                  "📨 Send Email\n(Outlook)\nto Patient", CARD_BG, WHITE, font_size=12)

    # Email content
    email_content = [
        "Email includes:",
        "  • Patient name & appointment date/time",
        "  • Clinic name & appointment category",
        "  • Risk score percentage",
        "  • Confirm/Reschedule call-to-action",
        "",
        "Payload: JSON with patient array + metadata",
        "Auth: Signed webhook URL (no OAuth needed)",
    ]
    add_bullet_text(slide, email_content, Inches(0.5), Inches(4.5), Inches(12), Inches(3), font_size=15)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 11: COST IMPACT & ROI
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Cost Impact & ROI", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    # Left column - Losses
    add_shape_box(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(4.5),
                  "WITHOUT Intervention\n━━━━━━━━━━━━━━━━\n\n"
                  "💰 Avg Revenue/Appt: $250\n\n"
                  "📉 Predicted No-Shows: ~35\n\n"
                  "❌ Estimated Loss: $8,750+\n\n"
                  "⚠️ Based on 594 upcoming\n     appointments in pipeline",
                  CARD_BG, WHITE, font_size=15)

    # Right column - Savings
    add_shape_box(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(4.5),
                  "WITH Risk-Based Intervention\n━━━━━━━━━━━━━━━━━━━━\n\n"
                  "🎯 High-Risk Callbacks: 40% recovery\n\n"
                  "📱 Medium-Risk SMS: 25% recovery\n\n"
                  "✅ Recoverable Revenue: ~$2,500+\n\n"
                  "📈 ROI: Positive from Day 1\n     (zero infrastructure cost)",
                  CARD_BG, WHITE, font_size=15)

    add_title_text(slide, "Zero marginal cost — all tools are free tier or existing Microsoft licenses",
                   Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
                   font_size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 12: TECH STACK & DEMO
    # ══════════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)

    add_title_text(slide, "Tech Stack & Live Demo", Inches(0.5), Inches(0.3), Inches(8), Inches(0.8),
                   font_size=32, color=TEAL)

    # Tech stack grid
    stack = [
        ("Data Layer", "Power BI Semantic Model, DAX, REST API"),
        ("Auth", "Azure AD Service Principal (OAuth2)"),
        ("ML Engine", "XGBoost (scikit-learn compatible)"),
        ("Frontend", "Streamlit (Python, deployed on Streamlit Cloud)"),
        ("AI/NLP", "Groq API — Llama 3.1 8B Instant"),
        ("Automation", "Power Automate (HTTP webhook trigger)"),
        ("Hosting", "Streamlit Cloud (free tier)"),
        ("Version Control", "GitHub (main branch, CI/CD auto-deploy)"),
    ]

    y = Inches(1.5)
    for label, value in stack:
        add_title_text(slide, f"{label}:", Inches(0.5), y, Inches(2.5), Inches(0.4),
                       font_size=14, bold=True, color=TEAL)
        add_title_text(slide, value, Inches(3.2), y, Inches(9), Inches(0.4),
                       font_size=14, bold=False, color=WHITE)
        y += Inches(0.55)

    # Demo links
    add_shape_box(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(1.2),
                  "🌐 Live Demo: azureenterpriseintelligence.streamlit.app    |    "
                  "📂 GitHub: github.com/i-am-anmol/azure_enterprise_intelligence",
                  CARD_BG, TEAL, font_size=14)

    # ══════════════════════════════════════════════════════════════════════════
    # SAVE
    # ══════════════════════════════════════════════════════════════════════════
    output_path = os.path.join(os.path.dirname(__file__), "Clinical_NoShow_Prediction.pptx")
    prs.save(output_path)
    print(f"✅ Presentation saved: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
